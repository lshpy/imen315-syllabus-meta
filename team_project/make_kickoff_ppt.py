"""1차 킥오프 발표용 PPT — 내용 중심 깊이 버전."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

ROOT = Path(__file__).parent
CHARTS = ROOT / "charts"
FORMULAS = ROOT / "formulas"
SHOTS = ROOT / "screenshots"

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

INK = RGBColor(0x0F, 0x14, 0x1F)
INK_2 = RGBColor(0x42, 0x4B, 0x5C)
MUTED = RGBColor(0x8B, 0x95, 0xA7)
DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)
PAPER = RGBColor(0xFA, 0xFB, 0xFC)
ACCENT = RGBColor(0x4F, 0x46, 0xE5)
ACCENT_2 = RGBColor(0xEC, 0x48, 0x99)
ACCENT_3 = RGBColor(0x10, 0xB9, 0x81)
ACCENT_4 = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x18, 0x1B, 0x29)
FONT = "맑은 고딕"


def slide(bg=WHITE):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid(); rect.fill.fore_color.rgb = bg
    rect.line.fill.background()
    return s


def text(s, x, y, w, h, content, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, line_spacing=1.25):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    lines = content.split("\n") if "\n" in content else [content]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.alignment = align; p.line_spacing = line_spacing
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold
            r.font.color.rgb = color; r.font.name = FONT


def line(s, x1, y1, x2, y2, color=DIVIDER, weight=1.5):
    l = s.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    l.line.color.rgb = color; l.line.width = Pt(weight)


def chip(s, x, y, w, h, label, color=ACCENT, text_color=WHITE, size=10):
    from pptx.enum.text import MSO_ANCHOR
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background(); sh.adjustments[0] = 0.5
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top = tf.margin_bottom = Inches(0.0)
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = label; p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = text_color; r.font.name = FONT


def card(s, x, y, w, h, fill=PAPER, border=DIVIDER, border_width=0.75, accent=None, shadow=True):
    """카드: 미세 그림자 + 옵션 위쪽 액센트 바"""
    if shadow:
        sh_shadow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x + 0.03), Inches(y + 0.04),
                                        Inches(w), Inches(h))
        sh_shadow.fill.solid(); sh_shadow.fill.fore_color.rgb = RGBColor(0xEC, 0xEE, 0xF2)
        sh_shadow.line.fill.background(); sh_shadow.adjustments[0] = 0.05
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border_width > 0:
        sh.line.color.rgb = border; sh.line.width = Pt(border_width)
    else:
        sh.line.fill.background()
    sh.adjustments[0] = 0.05
    if accent is not None:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x + 0.15), Inches(y),
                                  Inches(w - 0.3), Inches(0.08))
        bar.fill.solid(); bar.fill.fore_color.rgb = accent
        bar.line.fill.background()


def page_no(s, n):
    text(s, 12.3, 7.05, 0.8, 0.3, f"{n:02d}", size=11, color=MUTED, align=PP_ALIGN.RIGHT)


def header(s, kicker, title, n):
    text(s, 0.7, 0.5, 12, 0.4, kicker, size=12, bold=True, color=ACCENT)
    text(s, 0.7, 0.9, 12, 1.0, title, size=32, bold=True, color=INK)
    line(s, 0.7, 2.0, 12.6, 2.0)
    page_no(s, n)


# ============================================
# 1. 표지
# ============================================
def s_cover():
    s = slide(bg=DARK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.4), Inches(0.6), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_2; bar.line.fill.background()

    text(s, 0.7, 2.6, 12, 0.4, "IMEN315 인간공학 · 1차 킥오프 · 2026.05.11",
         size=12, color=MUTED)
    text(s, 0.7, 3.0, 12, 0.9, "강의계획서의 인지·행동 유도 메커니즘",
         size=42, bold=True, color=WHITE)
    text(s, 0.7, 3.85, 12, 0.7, "인간공학적 역(逆)분석", size=42, bold=True, color=WHITE)
    text(s, 0.7, 5.0, 12, 1.0,
         "수업에서 배운 이론을 바로 그 수업의 강의계획서에 적용해\n학생 행동을 정량 예측하고 실제 학생들에게서 검증한다.",
         size=14, color=MUTED, line_spacing=1.5)
    line(s, 0.7, 6.6, 12.6, 6.6, color=RGBColor(0x35, 0x3B, 0x4D), weight=1)
    text(s, 0.7, 6.75, 6, 0.4, "팀장 이승현 · 산업경영공학부", size=12, color=MUTED)
    text(s, 6.7, 6.75, 6, 0.4, "오형석 교수 · 2026 봄학기", size=12, color=MUTED, align=PP_ALIGN.RIGHT)


# ============================================
# 2. 우리가 답하려는 질문
# ============================================
def s_question():
    s = slide()
    header(s, "RESEARCH QUESTION", "우리가 답하려는 핵심 질문", 2)

    text(s, 0.7, 2.4, 12, 0.4, "[Q1]", size=12, bold=True, color=ACCENT)
    text(s, 0.7, 2.75, 12, 0.7,
         "강의계획서의 4가지 설계 요소가",
         size=24, bold=True, color=INK)
    text(s, 0.7, 3.4, 12, 0.7,
         "인간공학 이론이 예측한 학생 행동을 실제로 유발하는가?",
         size=24, bold=True, color=ACCENT)

    line(s, 0.7, 4.4, 12.6, 4.4)

    text(s, 0.7, 4.6, 12, 0.4, "[Q2]", size=12, bold=True, color=ACCENT_2)
    text(s, 0.7, 4.95, 12, 0.7,
         "이론적 개선안 (분산 평가 · 프로필 카드 · gain frame)이",
         size=24, bold=True, color=INK)
    text(s, 0.7, 5.6, 12, 0.7,
         "행동 패턴을 의도된 방향으로 변화시키는가?",
         size=24, bold=True, color=ACCENT_2)

    text(s, 0.7, 6.7, 12, 0.4,
         "→ 답 = 이론 적용 깊이 × 검증 견고성 → 30점 핵심 항목 직격",
         size=12, color=MUTED)


# ============================================
# 3. 강의계획서 4규칙 (수업 안내서 그대로)
# ============================================
def s_4_rules():
    s = slide()
    header(s, "INPUT", "강의계획서가 학생에게 강요하는 규칙 4가지", 3)

    rules = [
        ("①", "랜덤 출석 확인 (10%)",
         "예정일 외에도 임의로 출석 확인",
         "→ 매일 출석 강박 발생", ACCENT_3),
        ("②", "단일 기말고사 (40%)",
         "중간고사 폐지, 8주차 휴강, 6/17 1회",
         "→ 학기말 학습 시간 폭증 (벼락치기)", ACCENT_4),
        ("③", "팀장 직접 지명 (최대 2명)",
         "지원서 미제출 시 자동 배정",
         "→ 친분 기반 팀 구성", ACCENT),
        ("④", "차등 인센티브",
         "팀장 +10점, 기여 저조 시 감점",
         "→ 무임승차·최소 기여 수렴", ACCENT_2),
    ]
    for i, (num, title, detail, behavior, color) in enumerate(rules):
        x = 0.7 + (i % 2) * 6.15
        y = 2.3 + (i // 2) * 2.4
        card(s, x, y, 5.95, 2.15)
        text(s, x + 0.3, y + 0.25, 0.7, 0.5, num, size=20, bold=True, color=color)
        text(s, x + 0.85, y + 0.25, 5, 0.55, title, size=17, bold=True, color=INK)
        text(s, x + 0.3, y + 0.95, 5.4, 0.45, detail, size=12, color=MUTED)
        line(s, x + 0.3, y + 1.5, x + 5.65, y + 1.5)
        text(s, x + 0.3, y + 1.6, 5.4, 0.45, behavior, size=13, bold=True, color=color)


# ============================================
# 4. 우리가 정량 예측하는 학생 행동
# ============================================
def s_predicted_behavior():
    s = slide()
    header(s, "EXPECTED BEHAVIOR", "강의계획서가 유발하는 4가지 행동 패턴", 4)

    items = [
        ("출석 habituation", "랜덤 체크에 적응 → 매일 출석 → 후반부 이탈",
         "출석률 6.7% → 20% (3×)", ACCENT_3),
        ("학기말 벼락치기", "기억 활성화가 학기 초 내용에서 급락",
         "인출 시간 7.4초 → 2.2초 (3.4×)", ACCENT_4),
        ("친분 기반 satisficing", "후보 평가의 작업 기억 한계 초과",
         "N>8에서 평가 완성도 70%↓", ACCENT),
        ("최소 기여 수렴 (무임승차)", "모호한 loss frame 회피 → 책임 회피 강함",
         "모호한 감점 표현이 +10 가산보다 자원 의도 ↑", ACCENT_2),
    ]
    for i, (name, mech, pred, color) in enumerate(items):
        y = 2.3 + i * 1.13
        card(s, 0.7, y, 11.9, 1.0)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y), Inches(0.12), Inches(1.0))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        text(s, 1.0, y + 0.18, 5.5, 0.4, name, size=16, bold=True, color=INK)
        text(s, 1.0, y + 0.6, 7.0, 0.35, mech, size=14, color=INK_2)
        text(s, 8.5, y + 0.32, 4.1, 0.4, pred, size=13, bold=True, color=color)


# ============================================
# 5. 이론 개요 (4개 한 페이지)
# ============================================
def s_theories():
    s = slide()
    header(s, "FRAMEWORK", "수업에서 배운 이론 4개로 정확히 예측", 5)

    items = [
        ("효용 학습", "수업 6주차 · Skinner",
         "f_utility.png",
         "가변비율 강화 — 슬롯머신 원리", ACCENT_3),
        ("망각 곡선 + 활성화", "수업 5–7주차 · Ebbinghaus + Anderson",
         "f_memory.png",
         "분산 학습이 인출을 가속", ACCENT_4),
        ("작업 기억 + 만족화", "수업 12주차 · Miller(1956) + Simon(1956)",
         "f_wm.png",
         "한계 초과 시 직관·친분으로 후퇴", ACCENT),
        ("표현 효과 + 손실 회피", "Tversky & Kahneman (1981)",
         "f_framing.png",
         "동일 결과도 손실 표현이 더 강함", ACCENT_2),
    ]
    for i, (name, src, fimg, desc, color) in enumerate(items):
        x = 0.7 + (i % 2) * 6.15
        y = 2.3 + (i // 2) * 2.4
        card(s, x, y, 5.95, 2.15, accent=color)
        text(s, x + 0.35, y + 0.25, 5.4, 0.5, name, size=18, bold=True, color=INK)
        text(s, x + 0.35, y + 0.75, 5.4, 0.35, src, size=12, color=MUTED)
        fpath = FORMULAS / fimg
        if fpath.exists():
            s.shapes.add_picture(str(fpath), Inches(x + 0.35), Inches(y + 1.1), height=Inches(0.5))
        text(s, x + 0.35, y + 1.7, 5.4, 0.4, desc, size=12, color=color)
    text(s, 0.7, 7.0, 12, 0.4,
         "→ 다음 4장에서 각 이론을 자세히 풀어봅니다.",
         size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============================================
# 5-A1. 이론 ① 효용 학습 — 쉬운 설명
# ============================================
def s_theory_utility_intuition():
    s = slide()
    header(s, "THEORY 1 / 4 · 직관", "보상이 가끔, 무작위로 나올 때 행동이 가장 잘 굳는다", 6)

    # 한 줄 메시지
    text(s, 0.7, 2.3, 12, 0.4, "쉽게 말하면", size=11, bold=True, color=ACCENT_3)
    text(s, 0.7, 2.7, 12, 0.7,
         "사람은 보상을 받으면 그 행동을 또 한다.",
         size=18, bold=True, color=INK)
    text(s, 0.7, 3.4, 12, 0.7,
         "근데 보상이 \"언제 나올지 모르게\" 띄엄띄엄 나올 때",
         size=18, color=INK_2)
    text(s, 0.7, 4.1, 12, 0.7,
         "그 행동이 가장 강하게, 가장 오래 굳는다.",
         size=20, bold=True, color=ACCENT_3)

    # 일상 비유 3개
    text(s, 0.7, 5.0, 12, 0.4, "이런 거랑 똑같아요", size=11, bold=True, color=INK)

    examples = [
        ("🎰", "슬롯머신", "당첨이 가끔, 언제 나올지 모름.\n끊기 어려움."),
        ("📱", "인스타·SNS 알림", "좋아요·댓글이 무작위 시점에 옴.\n계속 확인하게 됨."),
        ("🎮", "가챠 게임", "확률 뽑기. 다음 판에 나올 수도\n있다는 기대 → 멈출 수 없음."),
    ]
    for i, (emoji, name, desc) in enumerate(examples):
        x = 0.7 + i * 4.1
        card(s, x, 5.55, 3.9, 1.55, accent=ACCENT_3)
        text(s, x + 0.25, 5.75, 1, 0.5, emoji, size=28)
        text(s, x + 1.15, 5.78, 2.7, 0.4, name, size=15, bold=True, color=ACCENT_3)
        text(s, x + 0.25, 6.3, 3.5, 0.85, desc, size=13, color=INK_2, line_spacing=1.35)


# ============================================
# 5-A2. 이론 ① 효용 학습 — 공식·적용
# ============================================
def s_theory_utility():
    s = slide()
    header(s, "THEORY 1 / 4 · 적용", "효용 학습 — 공식과 강의계획서 적용", 7)

    # 좌측: 공식 + 변수 설명
    text(s, 0.7, 2.3, 5.7, 0.35, "수업 6주차 · 가변비율 강화 (Variable-Ratio Reinforcement)",
         size=11, bold=True, color=ACCENT_3)

    card(s, 0.7, 2.75, 5.95, 1.4, fill=RGBColor(0xF7, 0xF8, 0xFB), border_width=0)
    text(s, 0.95, 2.85, 5.5, 0.35, "핵심 공식", size=11, bold=True, color=MUTED)
    fp = FORMULAS / "f_utility.png"
    if fp.exists():
        s.shapes.add_picture(str(fp), Inches(0.95), Inches(3.3), height=Inches(0.7))

    text(s, 0.7, 4.3, 5.7, 0.35, "변수 의미", size=11, bold=True, color=INK)
    vars_ = [
        ("U(n)", "n번째 시행 후 행동 효용 (출석할 가치)"),
        ("R(n)", "이번 시행의 보상 (출석 체크 ⇒ +1 / 결석 적발 ⇒ −1)"),
        ("α", "학습률 (≈0.15) — 새 경험을 얼마나 받아들이는가"),
    ]
    for i, (sym, meaning) in enumerate(vars_):
        y = 4.7 + i * 0.45
        text(s, 0.95, y, 1, 0.35, sym, size=12, bold=True, color=ACCENT_3)
        text(s, 1.85, y, 4.5, 0.35, meaning, size=13, color=INK_2)

    text(s, 0.7, 6.2, 5.7, 0.35, "비유", size=11, bold=True, color=INK)
    text(s, 0.95, 6.55, 5.5, 0.45,
         "슬롯머신 — 보상이 \"가끔\" 무작위로 주어질 때\n행동이 가장 강하게 굳어짐",
         size=13, color=INK_2, line_spacing=1.3)

    card(s, 6.85, 2.75, 5.75, 4.55, fill=RGBColor(0xEC, 0xFD, 0xF5), border=ACCENT_3)
    text(s, 7.1, 2.95, 5, 0.4, "강의계획서가 만드는 효과", size=15, bold=True, color=ACCENT_3)
    text(s, 7.1, 3.45, 5.3, 0.6,
         "\"임의로 출석 체크\"",
         size=18, bold=True, color=INK)
    text(s, 7.1, 4.2, 5.3, 1.5,
         "오늘 안 가면 잡힐지 학생은 모른다.\n슬롯머신과 같은 원리로 매일 출석한다.\n학기 내내 같은 행동이 이어진다.",
         size=14, color=INK_2, line_spacing=1.45)

    line(s, 7.1, 5.85, 12.4, 5.85, color=ACCENT_3)
    text(s, 7.1, 5.95, 5, 0.4, "시뮬레이션 예측", size=12, bold=True, color=ACCENT_3)
    text(s, 7.1, 6.4, 5.3, 0.4, "고정 일정: 출석률 6.7%", size=14, color=INK_2)
    text(s, 7.1, 6.85, 5.3, 0.4, "랜덤+비공개: 20% (3배)", size=18, bold=True, color=INK)


# ============================================
# 5-B1. 이론 ② 망각 곡선 — 쉬운 설명
# ============================================
def s_theory_forget_intuition():
    s = slide()
    header(s, "THEORY 2 / 4 · 직관", "한 번 본 건 잊히고, 자주 나눠서 보면 오래 남는다", 8)

    text(s, 0.7, 2.3, 12, 0.4, "쉽게 말하면", size=11, bold=True, color=ACCENT_4)
    text(s, 0.7, 2.7, 12, 0.7,
         "기억은 시간이 지나면 사라진다.",
         size=18, bold=True, color=INK)
    text(s, 0.7, 3.4, 12, 0.7,
         "근데 \"여러 번에 나눠서\" 학습하면 사라지는 속도가 느려진다.",
         size=18, color=INK_2)
    text(s, 0.7, 4.1, 12, 0.7,
         "한 번에 몰아 학습한 것은 시험장에서 안 떠오른다.",
         size=20, bold=True, color=ACCENT_4)

    text(s, 0.7, 5.0, 12, 0.4, "이런 경험 다들 있죠", size=11, bold=True, color=INK)

    examples = [
        ("📚", "벼락치기", "시험 전날 6시간 공부 →\n시험장에서 \"분명 봤는데...\""),
        ("🎵", "노래 가사", "한 번 들으면 잊지만\n10번 띄엄띄엄 들으면 외워짐."),
        ("📖", "어릴 때 배운 것", "오래전에 자주 배워서\n지금까지 안 잊혀짐."),
    ]
    for i, (emoji, name, desc) in enumerate(examples):
        x = 0.7 + i * 4.1
        card(s, x, 5.55, 3.9, 1.55, accent=ACCENT_4)
        text(s, x + 0.25, 5.75, 1, 0.5, emoji, size=28)
        text(s, x + 1.15, 5.78, 2.7, 0.4, name, size=15, bold=True, color=ACCENT_4)
        text(s, x + 0.25, 6.3, 3.5, 0.85, desc, size=13, color=INK_2, line_spacing=1.35)


# ============================================
# 5-B2. 이론 ② 망각 곡선 — 공식·적용
# ============================================
def s_theory_forget():
    s = slide()
    header(s, "THEORY 2 / 4 · 적용", "망각 곡선 — 공식과 강의계획서 적용", 9)

    text(s, 0.7, 2.3, 5.7, 0.35, "수업 5–7주차 · Base-level Activation (Anderson)",
         size=11, bold=True, color=ACCENT_4)

    card(s, 0.7, 2.75, 5.95, 1.4, fill=RGBColor(0xF7, 0xF8, 0xFB), border_width=0)
    text(s, 0.95, 2.85, 5.5, 0.35, "핵심 공식", size=11, bold=True, color=MUTED)
    fp = FORMULAS / "f_memory.png"
    if fp.exists():
        s.shapes.add_picture(str(fp), Inches(0.95), Inches(3.25), height=Inches(0.75))

    text(s, 0.7, 4.3, 5.7, 0.35, "변수 의미", size=11, bold=True, color=INK)
    vars_ = [
        ("B", "청크의 활성도 (높을수록 잘 떠오름)"),
        ("t_j", "j번째 학습 후 경과 시간 (일)"),
        ("d", "망각 속도 (≈0.5)"),
        ("T_retrieval", "시험장에서 그 청크 떠올리는 시간 (초)"),
    ]
    for i, (sym, meaning) in enumerate(vars_):
        y = 4.7 + i * 0.42
        text(s, 0.95, y, 1.4, 0.35, sym, size=12, bold=True, color=ACCENT_4)
        text(s, 2.15, y, 4.2, 0.35, meaning, size=13, color=INK_2)

    text(s, 0.7, 6.55, 5.7, 0.35,
         "핵심: 자주·분산해서 학습할수록 B가 커지고 T가 짧아짐",
         size=11, color=MUTED)

    card(s, 6.85, 2.75, 5.75, 4.55, fill=RGBColor(0xFE, 0xF3, 0xC7), border=ACCENT_4)
    text(s, 7.1, 2.95, 5, 0.4, "강의계획서가 만드는 효과", size=15, bold=True, color=ACCENT_4)
    text(s, 7.1, 3.45, 5.3, 0.6,
         "\"기말고사 1회 40%\"",
         size=18, bold=True, color=INK)
    text(s, 7.1, 4.2, 5.3, 1.5,
         "3월에 배운 걸 6월 시험에서 떠올려야 한다.\n80일이 지나 기억은 거의 사라진 상태.\n결국 시험 직전 벼락치기.",
         size=14, color=INK_2, line_spacing=1.45)

    line(s, 7.1, 5.85, 12.4, 5.85, color=ACCENT_4)
    text(s, 7.1, 5.95, 5, 0.4, "시뮬레이션 — 시험장 인출 시간", size=12, bold=True, color=ACCENT_4)
    text(s, 7.1, 6.4, 5.3, 0.4, "단일 기말: 7.4초", size=14, color=INK_2)
    text(s, 7.1, 6.85, 5.3, 0.4, "격주 퀴즈: 2.2초 (3.4배 빠름)", size=18, bold=True, color=INK)


# ============================================
# 5-C1. 이론 ③ 작업 기억 — 쉬운 설명
# ============================================
def s_theory_wm_intuition():
    s = slide()
    header(s, "THEORY 3 / 4 · 직관", "사람 머리는 한 번에 7개만 비교 가능, 그 이상은 대충 결정", 10)

    text(s, 0.7, 2.3, 12, 0.4, "쉽게 말하면", size=11, bold=True, color=ACCENT)
    text(s, 0.7, 2.7, 12, 0.7,
         "사람은 동시에 여러 정보를 못 비교한다.",
         size=18, bold=True, color=INK)
    text(s, 0.7, 3.4, 12, 0.7,
         "정보가 많아지면 \"최적\" 대신 \"그럭저럭 괜찮은\" 첫 번째 옵션에서 멈춘다.",
         size=18, color=INK_2)
    text(s, 0.7, 4.25, 12, 0.7,
         "이걸 만족화 (Satisficing) 라고 한다 — Simon이 노벨상 받은 개념.",
         size=20, bold=True, color=ACCENT)

    text(s, 0.7, 5.15, 12, 0.4, "일상에서 매일 일어나는 일", size=11, bold=True, color=INK)

    examples = [
        ("☕", "카페 메뉴", "메뉴 3개면 다 비교.\n메뉴 30개면 \"그냥 라떼\""),
        ("🛒", "쇼핑몰", "비슷한 상품 50개 →\n중간 가격 첫 번째 클릭"),
        ("👥", "팀 멤버 뽑기", "후보 15명? → 떠오르는\n친한 5명으로 결정"),
    ]
    for i, (emoji, name, desc) in enumerate(examples):
        x = 0.7 + i * 4.1
        card(s, x, 5.55, 3.9, 1.55, accent=ACCENT)
        text(s, x + 0.25, 5.75, 1, 0.5, emoji, size=28)
        text(s, x + 1.15, 5.78, 2.7, 0.4, name, size=15, bold=True, color=ACCENT)
        text(s, x + 0.25, 6.3, 3.5, 0.85, desc, size=13, color=INK_2, line_spacing=1.35)


# ============================================
# 5-C2. 이론 ③ 작업 기억 — 공식·적용
# ============================================
def s_theory_wm():
    s = slide()
    header(s, "THEORY 3 / 4 · 적용", "작업 기억 — 7±2와 강의계획서 적용", 11)

    text(s, 0.7, 2.3, 5.7, 0.35,
         "수업 12주차 · Bounded Rationality (제한된 합리성)",
         size=11, bold=True, color=ACCENT)

    card(s, 0.7, 2.75, 5.95, 1.8, fill=RGBColor(0xF7, 0xF8, 0xFB), border_width=0)
    text(s, 0.95, 2.9, 5.5, 0.4, "두 가지 원리", size=13, bold=True, color=MUTED)
    fp = FORMULAS / "f_wm.png"
    if fp.exists():
        s.shapes.add_picture(str(fp), Inches(0.95), Inches(3.3), height=Inches(0.55))
    text(s, 0.95, 4.0, 5.5, 0.5,
         "한계 초과 시 만족화(Satisficing)로 후퇴",
         size=15, bold=True, color=INK)

    text(s, 0.7, 4.85, 5.7, 0.4,
         "만족화란",
         size=14, bold=True, color=INK)
    text(s, 0.95, 5.25, 5.5, 1.2,
         "최적 선택을 포기하고\n첫 번째로 \"괜찮은\" 옵션에서 멈춘다.",
         size=15, color=INK_2, line_spacing=1.4)

    text(s, 0.7, 6.5, 5.7, 0.4, "비유 — 카페 메뉴", size=14, bold=True, color=INK)
    text(s, 0.95, 6.9, 5.5, 0.5,
         "메뉴 3개면 다 비교. 메뉴 30개면 그냥 라떼.",
         size=14, color=MUTED)

    card(s, 6.85, 2.75, 5.75, 4.55, fill=RGBColor(0xEE, 0xEA, 0xFE), border=ACCENT)
    text(s, 7.1, 2.95, 5, 0.4, "강의계획서가 만드는 효과", size=15, bold=True, color=ACCENT)
    text(s, 7.1, 3.45, 5.3, 0.6,
         "\"팀장이 2명 직접 지명\"",
         size=20, bold=True, color=INK)
    text(s, 7.1, 4.2, 5.3, 1.5,
         "후보가 많으면 머리에 다 들어오지 않는다.\n떠오르는 친한 사람으로 결정한다.\n친분 기반 팀, 역량 균형이 깨진다.",
         size=15, color=INK_2, line_spacing=1.5)

    line(s, 7.1, 5.95, 12.4, 5.95, color=ACCENT)
    text(s, 7.1, 6.05, 5, 0.4, "시뮬레이션 — 팀 만족도", size=14, bold=True, color=ACCENT)
    text(s, 7.1, 6.5, 5.3, 0.4, "정보 없이 15명 중: 70% 이하", size=15, color=INK_2)
    text(s, 7.1, 6.95, 5.3, 0.4, "프로필 카드 제공: 92%", size=18, bold=True, color=INK)
    text(s, 7.1, 6.85, 5.3, 0.4, "프로필 카드 제공: 92%", size=16, bold=True, color=INK)


# ============================================
# 5-D1. 이론 ④ Framing — 쉬운 설명
# ============================================
def s_theory_framing_intuition():
    s = slide()
    header(s, "THEORY 4 / 4 · 직관", "같은 결과도 \"이득\" vs \"손실\" 표현에 따라 반응이 다르다", 12)

    text(s, 0.7, 2.3, 12, 0.4, "쉽게 말하면", size=11, bold=True, color=ACCENT_2)
    text(s, 0.7, 2.7, 12, 0.7,
         "사람은 \"잃는 것\"을 \"얻는 것\"보다 약 2배 더 강하게 느낀다.",
         size=18, bold=True, color=INK)
    text(s, 0.7, 3.4, 12, 0.7,
         "그래서 같은 -10점도 \"가산점 없음\"보다 \"감점\"이라 하면 훨씬 무서움.",
         size=18, color=INK_2)
    text(s, 0.7, 4.25, 12, 0.7,
         "→ 똑같은 정책도 표현 방식만 바꾸면 행동이 달라진다.",
         size=20, bold=True, color=ACCENT_2)

    text(s, 0.7, 5.15, 12, 0.4, "이런 차이 다들 느껴봤을 거예요", size=11, bold=True, color=INK)

    examples = [
        ("💰", "할인 vs 추가요금", "\"5% 할인\" 좋아하지만\n\"5% 추가\"는 짜증."),
        ("🏥", "수술 성공률", "\"90% 성공\" → 받음.\n\"10% 사망\" → 거부."),
        ("🚇", "교통카드", "\"적립 0원\"보다\n\"수수료 100원\"이 화남."),
    ]
    for i, (emoji, name, desc) in enumerate(examples):
        x = 0.7 + i * 4.1
        card(s, x, 5.55, 3.9, 1.55, accent=ACCENT_2)
        text(s, x + 0.25, 5.75, 1, 0.5, emoji, size=28)
        text(s, x + 1.15, 5.78, 2.7, 0.4, name, size=15, bold=True, color=ACCENT_2)
        text(s, x + 0.25, 6.3, 3.5, 0.85, desc, size=13, color=INK_2, line_spacing=1.35)


# ============================================
# 5-D2. 이론 ④ Framing — 공식·적용
# ============================================
def s_theory_framing():
    s = slide()
    header(s, "THEORY 4 / 4 · 적용", "표현 방식 효과 — 공식과 강의계획서 적용", 13)

    text(s, 0.7, 2.3, 5.7, 0.35,
         "Prospect Theory (1979) · Framing of Decisions (1981)",
         size=11, bold=True, color=ACCENT_2)

    card(s, 0.7, 2.75, 5.95, 1.7, fill=PAPER)
    text(s, 0.95, 2.95, 5.5, 0.4, "핵심 원리", size=11, bold=True, color=MUTED)
    text(s, 0.95, 3.35, 5.5, 0.5,
         "동일한 결과도 \"이득\"으로 표현 vs \"손실\"로 표현",
         size=13, bold=True, color=INK)
    text(s, 0.95, 3.85, 5.5, 0.5,
         "→ 사람의 반응이 다르다",
         size=13, bold=True, color=INK)

    text(s, 0.7, 4.7, 5.7, 0.4, "수치적 비대칭", size=11, bold=True, color=INK)
    fp = FORMULAS / "f_framing.png"
    if fp.exists():
        s.shapes.add_picture(str(fp), Inches(0.95), Inches(5.05), height=Inches(0.55))
    text(s, 0.95, 5.7, 5.5, 0.5,
         "= 같은 −10원 손실이 +10원 이득보다 약 2배 강함",
         size=13, color=INK_2)

    text(s, 0.7, 6.45, 5.7, 0.35, "고전 실험 — Asian disease problem", size=11, bold=True, color=INK)
    text(s, 0.95, 6.8, 5.5, 0.4,
         "동일 정책 → 표현만 바꿔도 선택률 72% vs 22%",
         size=11, color=MUTED)

    card(s, 6.85, 2.75, 5.75, 4.55, fill=RGBColor(0xFD, 0xE8, 0xF5), border=ACCENT_2)
    text(s, 7.1, 2.95, 5, 0.4, "강의계획서가 만드는 효과", size=15, bold=True, color=ACCENT_2)
    text(s, 7.1, 3.45, 5.3, 0.5,
         "\"+10 가산\" vs \"감점\"",
         size=18, bold=True, color=INK)
    text(s, 7.1, 4.1, 5.3, 1.5,
         "+10은 명확한 보상.\n감점은 기준이 모호한 위협.\n학생은 모호한 위협에 더 강하게 회피한다.",
         size=14, color=INK_2, line_spacing=1.45)

    line(s, 7.1, 5.85, 12.4, 5.85, color=ACCENT_2)
    text(s, 7.1, 5.95, 5, 0.4, "예상 결과", size=12, bold=True, color=ACCENT_2)
    text(s, 7.1, 6.4, 5.3, 0.4, "+10 보고 자원: 적당함", size=14, color=INK_2)
    text(s, 7.1, 6.85, 5.3, 0.4, "감점 피하려 자원: 약 2배", size=18, bold=True, color=INK)


# ============================================
# 6. 우리의 핵심 가설 4개
# ============================================
def s_hypotheses():
    s = slide()
    header(s, "HYPOTHESES", "우리가 데이터로 검증할 4가지 가설", 6)

    h = [
        ("H1", "표현 방식 효과", "ACCENT_2",
         "모호한 감점 표현이 명확한 +10 가산보다 회피·자원 의도가 더 강하다.",
         "감점 그룹 응답이 +10 그룹보다 약 2배 더 강한 의도 (Cohen's d ≈ 0.5+)"),
        ("H2", "그럭저럭 만족", "ACCENT",
         "후보 N이 7명을 넘어가면 팀 구성 만족도가 급격히 떨어진다.",
         "N=3·7 그룹은 만족도 5~6점, N=15 그룹은 4점 이하로 급감"),
        ("H3", "효용 학습", "ACCENT_3",
         "출석 정책이 무작위+비공개일수록 결석 의도가 가장 낮다.",
         "결석 의도: 고정 매주월 > 랜덤 빈도공개 > 랜덤 비공개 순서로 감소"),
        ("H4", "망각 곡선", "ACCENT_4",
         "단일 기말 체제는 격주 퀴즈 체제보다 학습시간 분산이 훨씬 크다.",
         "단일 기말: 8주차 5h, 15주차 15h (3배). 격주 퀴즈: 8주차·15주차 ≈ 9h"),
    ]
    color_map = {"ACCENT": ACCENT, "ACCENT_2": ACCENT_2, "ACCENT_3": ACCENT_3, "ACCENT_4": ACCENT_4}
    for i, (code, theory, color_key, body, expected) in enumerate(h):
        color = color_map[color_key]
        y = 2.3 + i * 1.18
        card(s, 0.7, y, 11.9, 1.05)
        chip(s, 0.9, y + 0.3, 0.85, 0.45, code, color=color, size=11)
        text(s, 2.0, y + 0.18, 4.0, 0.35, theory, size=11, bold=True, color=color)
        text(s, 2.0, y + 0.48, 10.5, 0.4, body, size=12, bold=True, color=INK)
        text(s, 2.0, y + 0.78, 10.5, 0.35, "기대 결과 — " + expected, size=12, color=INK_2)


# ============================================
# T1 개요 — 줄글 설명
# ============================================
def s_t1_overview():
    s = slide()
    header(s, "T1 · OVERVIEW", "T1 설문 실험 — 정확히 무엇을 어떻게 하는가", 14)

    # 본문 1: 큰 한 줄
    text(s, 0.7, 2.3, 12, 0.5,
         "한 응답자가 5분 안에 4가지 가상 상황을 마주합니다.",
         size=18, bold=True, color=INK)
    text(s, 0.7, 2.85, 12, 0.5,
         "그 응답을 바탕으로 우리는 인간공학 이론 4개의 예측이 맞는지 검증합니다.",
         size=16, color=INK_2)

    line(s, 0.7, 3.65, 12.6, 3.65)

    # 본문 2: 어떻게 진행되나
    text(s, 0.7, 3.85, 12, 0.4, "응답자 입장에서 진행되는 흐름", size=11, bold=True, color=ACCENT)

    flows = [
        ("①", "도입", "닉네임 + 학년 입력 (5초)"),
        ("②", "출석 정책 비교", "비 오는 화요일에 출석할지를 3 정책에서 답함 (within-subjects)"),
        ("③", "팀장 인센티브", "강의계획서의 gain·loss 표현 중 하나를 무작위로 보고 자원 의도 응답 (between)"),
        ("④", "팀원 직접 지명", "후보 N명 중 2명 지명 — 친분 vs 역량 트레이드오프 (between)"),
        ("⑤", "평가 체계 시뮬", "두 평가 체제에서 본인 학습시간 분포 추정 (within)"),
        ("⑥", "결과 분석", "Streamlit은 16 학습 MBTI 자동 진단 + 즉시 피드백"),
    ]
    for i, (num, name, desc) in enumerate(flows):
        y = 4.3 + i * 0.42
        text(s, 0.95, y, 0.4, 0.35, num, size=14, bold=True, color=ACCENT)
        text(s, 1.45, y, 2.2, 0.35, name, size=12, bold=True, color=INK)
        text(s, 3.7, y, 8.9, 0.35, desc, size=13, color=INK_2)

    text(s, 0.7, 7.0, 12, 0.4,
         "가장 중요한 것: 한 명의 응답이 H1·H2·H3·H4 4가지 가설을 동시에 검증하는 데이터가 됨",
         size=11, bold=True, color=ACCENT_2, align=PP_ALIGN.CENTER)


# ============================================
# T1 데이터 흐름 — 줄글
# ============================================
def s_t1_dataflow():
    s = slide()
    header(s, "T1 · DATA FLOW", "응답 → 데이터 → 검증 — 한눈에 흐름", 15)

    # 1단계 박스들
    steps = [
        ("1️⃣ 수집", "Streamlit + Google Forms\n동시 배포 (단톡, 에타, 페북)", ACCENT_3),
        ("2️⃣ 누적", "Google Sheets\n실시간 자동 적재", ACCENT),
        ("3️⃣ 정제", "조작 점검 통과자만 필터\n결측치 처리", ACCENT_4),
        ("4️⃣ 검정", "Python (pandas, scipy, pingouin)\nt-test, ANOVA, F-test", ACCENT_2),
    ]
    for i, (title, desc, color) in enumerate(steps):
        x = 0.7 + i * 3.18
        card(s, x, 2.5, 2.95, 2.1, fill=PAPER, border=color)
        text(s, x + 0.2, 2.65, 2.6, 0.5, title, size=15, bold=True, color=color)
        text(s, x + 0.2, 3.2, 2.6, 1.3, desc, size=13, color=INK_2, line_spacing=1.4)
        if i < 3:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        Inches(x + 2.95), Inches(3.4),
                                        Inches(0.23), Inches(0.3))
            arrow.fill.solid(); arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    # 검증 결과 예시 박스
    line(s, 0.7, 4.95, 12.6, 4.95)
    text(s, 0.7, 5.1, 12, 0.4, "최종 산출 — 보고서·발표에 그대로 들어가는 결과 형식",
         size=12, bold=True, color=ACCENT)

    outputs = [
        ("📊", "통계 결과표",
         "각 가설별 t값/F값/p값/효과크기 표"),
        ("📈", "예측 vs 실측 그래프",
         "시뮬 예측치와 학생 응답 평균을 같은 축에"),
        ("🧬", "16 유형 분포",
         "수강생 중 어느 학습 MBTI가 많은지 분포도"),
    ]
    for i, (emoji, name, desc) in enumerate(outputs):
        x = 0.7 + i * 4.1
        text(s, x, 5.65, 0.5, 0.4, emoji, size=20)
        text(s, x + 0.6, 5.65, 3.4, 0.4, name, size=13, bold=True, color=INK)
        text(s, x + 0.6, 6.05, 3.4, 0.7, desc, size=13, color=INK_2, line_spacing=1.3)

    text(s, 0.7, 7.0, 12, 0.4,
         "끝나면 N=60 응답으로 4 가설 검정 + 시뮬 예측 정합성 확인 + 도구 A/B 응답률 비교",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================
# 7. T1 시나리오 1: 출석
# ============================================
def s_t1_attend():
    s = slide()
    header(s, "T1 · 시나리오 ①", "출석 정책에 따른 결석 의도 (within-subjects)", 7)

    text(s, 0.7, 2.35, 12, 0.4,
         "한 응답자가 3가지 가상 상황을 모두 평가 → 동일인 비교로 검정력 ↑",
         size=12, color=MUTED)

    scenarios = [
        ("A", "고정 — 매주 월요일 체크", "고정 일정", "출석률 6.7% (예측)", ACCENT_3),
        ("B", "랜덤 + 횟수 공지 (30회 중 7회)", "가변·정보 일부", "출석률 20% (예측)", ACCENT),
        ("C", "랜덤 + 횟수 비공개", "가변·정보 무", "출석률 20% (가변강화 최대)", ACCENT_2),
    ]
    for i, (code, title, label, pred, color) in enumerate(scenarios):
        y = 2.85 + i * 1.2
        chip(s, 0.7, y + 0.3, 0.6, 0.45, code, color=color)
        text(s, 1.5, y + 0.18, 6.7, 0.4, title, size=14, bold=True, color=INK)
        text(s, 1.5, y + 0.6, 6.7, 0.35, "→ 응답: \"이 정책에서 비 오는 날 결석할 의도\" (1~7)",
             size=12, color=INK_2)
        text(s, 8.5, y + 0.32, 4.2, 0.4, pred, size=11, bold=True, color=color, align=PP_ALIGN.RIGHT)
        if i < 2:
            line(s, 0.7, y + 1.1, 12.6, y + 1.1)

    text(s, 0.7, 6.95, 12, 0.4,
         "측정: 결석 의도 1~7 점수 + 비 오는 날 갈 가능성 0~100% · 가설 H3 검정",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================
# 8. T1 시나리오 2: 인센티브 표현 (Framing) — 강의계획서 그대로
# ============================================
def s_t1_framing():
    s = slide()
    header(s, "T1 · 시나리오 ②", "팀장 인센티브 표현 효과 — 강의계획서 그대로", 16)

    text(s, 0.7, 2.4, 12, 0.4,
         "강의계획서의 두 표현을 그대로 분리해 보여줌 — gain은 명확, loss는 모호 (비대칭 자체가 발견)",
         size=12, color=MUTED)

    # 좌: GAIN frame (명확)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(3.0), Inches(0.1), Inches(3.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_3; bar.line.fill.background()
    card(s, 0.85, 3.0, 5.8, 3.5, fill=RGBColor(0xEC, 0xFD, 0xF5), border=ACCENT_3)
    text(s, 1.05, 3.2, 5, 0.4, "조건 A · GAIN frame (명확)", size=15, bold=True, color=ACCENT_3)
    text(s, 1.05, 3.65, 5.5, 0.85,
         "\"팀장 자원자에게\n+10점 가산\"",
         size=18, bold=True, color=INK, line_spacing=1.25)
    text(s, 1.05, 4.95, 5.5, 0.4, "표현 특성", size=11, bold=True, color=MUTED)
    text(s, 1.05, 5.3, 5.5, 0.4, "+10이라는 명확한 수치", size=14, color=INK_2)
    text(s, 1.05, 5.7, 5.5, 0.4, "능동 보상 · 긍정 톤", size=14, color=INK_2)
    text(s, 1.05, 6.15, 5.5, 0.5, "예상: 자원 의도 약함", size=12, bold=True, color=ACCENT_3)

    # 우: LOSS frame (모호)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(3.0), Inches(0.1), Inches(3.5))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_2; bar.line.fill.background()
    card(s, 7.0, 3.0, 5.6, 3.5, fill=RGBColor(0xFD, 0xE8, 0xF5), border=ACCENT_2)
    text(s, 7.2, 3.2, 5, 0.4, "조건 B · LOSS frame (모호)", size=15, bold=True, color=ACCENT_2)
    text(s, 7.2, 3.65, 5.3, 0.85,
         "\"기여도 저조 시 감점\n(현저히 낮다고 판단 시)\"",
         size=16, bold=True, color=INK, line_spacing=1.25)
    text(s, 7.2, 4.95, 5.3, 0.4, "표현 특성", size=11, bold=True, color=MUTED)
    text(s, 7.2, 5.3, 5.3, 0.4, "감점 폭 미명시 → 모호", size=14, color=INK_2)
    text(s, 7.2, 5.7, 5.3, 0.4, "위협 · 회피 · 손실 톤", size=14, color=INK_2)
    text(s, 7.2, 6.15, 5.3, 0.5, "예상: 회피 의도 강함 (≈2×)", size=12, bold=True, color=ACCENT_2)

    text(s, 0.7, 6.85, 12, 0.4,
         "발견 포인트: 강의계획서가 gain은 명확하게, loss는 모호하게 설계 → 학생 회피 행동 강화",
         size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============================================
# 9. T1 시나리오 3: 팀 구성 (실제 데이터 수집)
# ============================================
def s_t1_team():
    s = slide()
    header(s, "T1 · 시나리오 ③", "팀원 직접 지명 — 실제 행동 측정", 9)

    text(s, 0.7, 2.4, 12, 0.5,
         "후보 N명에서 정확히 2명 지명 (강의계획서 규칙 그대로) — 결정 시간·만족도·친분 편향 측정",
         size=12, color=MUTED)

    # 3 그룹 박스
    groups = [
        ("N = 3", "단순 비교 가능", "WM 한계 내", ACCENT_3, 0.7),
        ("N = 7", "WM 임계점", "처리 한계 직전", ACCENT_4, 4.85),
        ("N = 15", "한계 초과", "Satisficing 발현", ACCENT_2, 9.0),
    ]
    for label, desc, status, color, x in groups:
        card(s, x, 3.0, 3.83, 1.7)
        text(s, x + 0.2, 3.15, 3.5, 0.5, label, size=22, bold=True, color=color)
        text(s, x + 0.2, 3.7, 3.5, 0.4, desc, size=13, color=INK_2)
        text(s, x + 0.2, 4.15, 3.5, 0.4, status, size=10, bold=True, color=color)

    # 측정 지표
    text(s, 0.7, 5.0, 12, 0.4, "자동 측정 지표 (응답 페이지 timestamp 활용)", size=12, bold=True, color=ACCENT)

    metrics = [
        ("친분 편향 점수", "선택된 2명의 친밀도 / 최댓값 (0~1)"),
        ("역량 점수", "선택된 2명의 역량 / 최댓값 (0~1)"),
        ("결정 시간", "페이지 진입 → 제출까지 (초)"),
        ("만족도", "1~7 Likert"),
    ]
    for i, (m, d) in enumerate(metrics):
        y = 5.5 + (i // 2) * 0.55
        x = 0.95 + (i % 2) * 6.15
        text(s, x, y, 6, 0.35, "▸ " + m, size=11, bold=True, color=INK)
        text(s, x, y + 0.27, 6, 0.3, d, size=11, color=MUTED)

    text(s, 0.7, 6.95, 12, 0.4,
         "후보 풀 설계: 친분 ★★★★★(역량 중) ↔ 친분 ★(역량 최고) — 트레이드오프 강제",
         size=11, color=ACCENT_2, align=PP_ALIGN.CENTER)


# ============================================
# 10. T1 시나리오 4: 평가 체계
# ============================================
def s_t1_eval():
    s = slide()
    header(s, "T1 · 시나리오 ④", "평가 체계 → 학습시간 분포 (within-subjects)", 10)

    text(s, 0.7, 2.4, 12, 0.4,
         "한 응답자가 두 평가 체계 시나리오에서 본인의 학습 시간을 시점별로 예상",
         size=12, color=MUTED)

    # 두 시나리오
    card(s, 0.7, 2.95, 5.95, 3.6)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.95), Inches(0.1), Inches(3.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_4; bar.line.fill.background()
    text(s, 0.95, 3.1, 5, 0.5, "체제 A · 단일 기말 (40%)", size=14, bold=True, color=ACCENT_4)
    text(s, 0.95, 3.6, 5.5, 0.4, "응답: 8주차 · 15주차 학습/주 (시간)", size=12, color=MUTED)
    text(s, 0.95, 4.1, 5.5, 0.4, "예상 분포", size=11, bold=True, color=INK)
    text(s, 0.95, 4.5, 5.5, 0.4, "8주차: 적음 · 15주차: 폭증", size=14, color=INK_2)
    line(s, 0.95, 5.0, 6.4, 5.0)
    text(s, 0.95, 5.1, 5.5, 0.4, "이론 예측", size=11, bold=True, color=ACCENT_4)
    text(s, 0.95, 5.45, 5.5, 0.7, "B 활성화 ≈ −2.0\n인출 시간 7.4초", size=14, bold=True, color=INK)

    card(s, 6.85, 2.95, 5.75, 3.6)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.85), Inches(2.95), Inches(0.1), Inches(3.6))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_3; bar.line.fill.background()
    text(s, 7.1, 3.1, 5, 0.5, "체제 B · 격주 퀴즈 6회 + 기말", size=14, bold=True, color=ACCENT_3)
    text(s, 7.1, 3.6, 5, 0.4, "응답: 8주차 · 15주차 학습/주 (시간)", size=12, color=MUTED)
    text(s, 7.1, 4.1, 5, 0.4, "예상 분포", size=11, bold=True, color=INK)
    text(s, 7.1, 4.5, 5, 0.4, "8주차 ≈ 15주차 (분산 학습)", size=14, color=INK_2)
    line(s, 7.1, 5.0, 12.5, 5.0)
    text(s, 7.1, 5.1, 5, 0.4, "이론 예측", size=11, bold=True, color=ACCENT_3)
    text(s, 7.1, 5.45, 5, 0.7, "B 활성화 ≈ −0.8\n인출 시간 2.2초 (3.4× 빠름)", size=14, bold=True, color=INK)

    text(s, 0.7, 6.85, 12, 0.4,
         "측정: 학습시간 분포 분산 비교 (F-test) · 이론 예측치 vs 학생 자기보고 정합성",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================
# 11. T1 배포 — 도구 3종 A/B/C
# ============================================
def s_t1_distribution():
    s = slide()
    header(s, "T1 · DISTRIBUTION", "동일 설문, 3가지 도구 동시 배포 — 인터페이스 효과 측정", 11)

    text(s, 0.7, 2.4, 12, 0.4,
         "도구 자체의 응답률 차이 = 인간공학적 UI 설계 효과의 실증 데이터",
         size=12, color=MUTED)

    tools = [
        ("A", "Streamlit 게임형", "학습 MBTI 16유형 분석",
         "Fitts ✓ 작업기억 ✓ 즉시 피드백 ✓",
         "예상 응답률 ~45%", ACCENT_3),
        ("B", "Google Forms 대조군", "표준 설문 폼",
         "참조 베이스라인",
         "예상 응답률 ~30%", INK_2),
        ("C", "Google Sheets 응답 시트", "실시간 누적 + 분석 대시보드",
         "응답 자동 적재 · 차트 자동 갱신",
         "Admin: ?admin=lshpy2026", ACCENT),
    ]
    for i, (code, title, desc, principle, expect, color) in enumerate(tools):
        x = 0.7 + i * 4.1
        card(s, x, 2.95, 3.9, 4.0)
        chip(s, x + 0.2, 3.1, 0.5, 0.4, code, color=color)
        text(s, x + 0.95, 3.1, 3, 0.4, title, size=14, bold=True, color=INK)
        text(s, x + 0.2, 3.7, 3.6, 0.4, desc, size=13, color=INK_2)
        line(s, x + 0.2, 4.2, x + 3.7, 4.2)
        text(s, x + 0.2, 4.3, 3.6, 0.4, "특징", size=11, bold=True, color=MUTED)
        text(s, x + 0.2, 4.65, 3.6, 0.9, principle, size=13, color=INK_2)
        line(s, x + 0.2, 5.7, x + 3.7, 5.7)
        text(s, x + 0.2, 5.8, 3.6, 0.4, "결과 가설", size=11, bold=True, color=MUTED)
        text(s, x + 0.2, 6.15, 3.6, 0.4, expect, size=12, bold=True, color=color)


# ============================================
# 12. T2: 컴퓨터 시뮬 (쉽게 풀어쓰기)
# ============================================
def s_t2():
    s = slide()
    header(s, "T2 · 컴퓨터 시뮬레이션", "가상 학생 1000명을 만들어 \"약한 학생도 보호되는지\" 본다", 22)

    # 한 줄 메시지
    text(s, 0.7, 2.3, 12, 0.45,
         "왜 1000명이나? 평균만 보면 가려지는 진실이 있다.",
         size=15, bold=True, color=INK)
    text(s, 0.7, 2.8, 12, 0.45,
         "예: \"평균 출석률 80%\"라고 해도, 학습률 낮은 학생은 30%일 수 있음.",
         size=13, color=INK_2)

    # 큰 비교
    line(s, 0.7, 3.6, 12.6, 3.6)
    text(s, 0.7, 3.75, 12, 0.4, "개인과제 → 팀과제 확장", size=11, bold=True, color=ACCENT)

    card(s, 0.7, 4.2, 5.95, 2.6, fill=PAPER)
    text(s, 0.95, 4.35, 5.5, 0.4, "개인과제 (현재)", size=12, bold=True, color=MUTED)
    text(s, 0.95, 4.8, 5.5, 0.5, "학생 1명 평균만", size=18, bold=True, color=INK)
    text(s, 0.95, 5.4, 5.5, 1.3,
         "모든 학생이 똑같이 학습률 0.15,\n작업 기억 7개라고 가정.\n→ \"평균적 학생\"만 본 결과.",
         size=13, color=INK_2, line_spacing=1.4)

    card(s, 6.85, 4.2, 5.75, 2.6, fill=RGBColor(0xEE, 0xEA, 0xFE), border=ACCENT)
    text(s, 7.1, 4.35, 5.3, 0.4, "팀과제 (확장)", size=12, bold=True, color=ACCENT)
    text(s, 7.1, 4.8, 5.3, 0.5, "가상 학생 1000명, 사람마다 다르게", size=15, bold=True, color=INK)
    text(s, 7.1, 5.4, 5.3, 1.3,
         "학습률·기억력·작업 기억이 모두 다른\n1000명을 가상으로 만들어 시뮬.\n→ 평균뿐 아니라 하위 25%도 본다.",
         size=13, color=INK_2, line_spacing=1.4)

    text(s, 0.7, 7.0, 12, 0.4,
         "분석 질문: 개선안이 약한 학생도 도와주는가",
         size=12, bold=True, color=ACCENT_2, align=PP_ALIGN.CENTER)


# ============================================
# 13. T3: 친분 그래프 (쉽게)
# ============================================
def s_t3():
    s = slide()
    header(s, "T3 · 친분 그래프 시뮬", "친구 관계를 그림으로 그려서 \"팀장이 누구를 뽑는지\" 시뮬", 23)

    # 한 줄 메시지
    text(s, 0.7, 2.3, 12, 0.45,
         "사람 모집 없이도 친분 휴리스틱을 검증할 수 있다.",
         size=15, bold=True, color=INK)
    text(s, 0.7, 2.8, 12, 0.45,
         "가상 학생 30명에 친한 정도를 다 설정 → 팀장이 5가지 방식으로 뽑게 함 → 그림으로 비교.",
         size=14, color=INK_2)

    line(s, 0.7, 3.6, 12.6, 3.6)
    text(s, 0.7, 3.75, 12, 0.4, "5가지 팀장 선택 방식 (그림 5장 산출)", size=11, bold=True, color=ACCENT)

    conditions = [
        ("①", "무작위로 뽑기", "그냥 랜덤 5명", "친분 편향 0 → 비교 기준", ACCENT_3),
        ("②", "친한 사람만", "친밀도 상위 5명", "친분 100%, 역량 낮음 (현실)", ACCENT_2),
        ("③", "잘하는 사람만", "역량 상위 5명", "역량 100%, 친분 0 (이상)", ACCENT),
        ("④", "머리에 떠오르는 사람", "친한 7명 떠올리고 그 중 5명", "현실의 satisficing", ACCENT_4),
        ("⑤", "프로필 카드 보고", "정보 본 후 다양성 고려", "균형 잡힌 팀", INK_2),
    ]
    for i, (code, name, alg, expect, color) in enumerate(conditions):
        y = 4.25 + i * 0.52
        text(s, 0.95, y + 0.05, 0.4, 0.35, code, size=14, bold=True, color=color)
        text(s, 1.5, y, 3.5, 0.4, name, size=12, bold=True, color=INK)
        text(s, 5.2, y, 4.0, 0.4, alg, size=12, color=INK_2)
        text(s, 9.4, y, 3.3, 0.4, expect, size=10, bold=True, color=color)

    text(s, 0.7, 7.0, 12, 0.4,
         "결과: 5장의 네트워크 그래프 → 발표 슬라이드 임팩트 최고",
         size=12, bold=True, color=ACCENT_2, align=PP_ALIGN.CENTER)


# ============================================
# NEW: Streamlit 데모 — 2장 분리 (크게 보이게)
# ============================================
def s_streamlit_demo_1():
    s = slide()
    header(s, "설문 화면 1", "응답자 화면 — 인트로 / 출석 시나리오", 21)

    path1 = SHOTS / "01_intro.png"
    if path1.exists():
        s.shapes.add_picture(str(path1), Inches(0.7), Inches(2.4), width=Inches(6.1))
    text(s, 0.7, 6.6, 6.1, 0.35, "인트로 화면",
         size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, 0.7, 6.95, 6.1, 0.3, "16 유형 미리 보여주고 닉네임 입력 후 시작",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)

    path2 = SHOTS / "02_attendance.png"
    if path2.exists():
        s.shapes.add_picture(str(path2), Inches(6.95), Inches(2.4), width=Inches(5.95))
    text(s, 6.95, 6.6, 5.95, 0.35, "출석 정책 비교 화면",
         size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, 6.95, 6.95, 5.95, 0.3, "3개 정책에서 각각 결석 의도를 골라요",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)


def s_streamlit_demo_2():
    s = slide()
    header(s, "설문 화면 2", "응답자 화면 — 팀장 인센티브 / 팀원 지명", 22)

    path1 = SHOTS / "03_framing.png"
    if path1.exists():
        s.shapes.add_picture(str(path1), Inches(0.7), Inches(2.4), width=Inches(6.1))
    text(s, 0.7, 6.6, 6.1, 0.35, "팀장 인센티브 화면",
         size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, 0.7, 6.95, 6.1, 0.3, "응답자 절반은 +10 가산, 절반은 모호한 감점 표현",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)

    path2 = SHOTS / "04_team.png"
    if path2.exists():
        s.shapes.add_picture(str(path2), Inches(6.95), Inches(2.4), width=Inches(5.95))
    text(s, 6.95, 6.6, 5.95, 0.35, "팀원 직접 지명 화면",
         size=12, bold=True, color=INK, align=PP_ALIGN.CENTER)
    text(s, 6.95, 6.95, 5.95, 0.3, "후보 N명 중 2명 지명 (강의계획서 규칙)",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)


# 호환: 기존 함수명 유지
def s_streamlit_demo():
    s_streamlit_demo_1()
    s_streamlit_demo_2()


# ============================================
# NEW: T2 시뮬 입력 예시 — 가상 학생 1명 자세히
# ============================================
def s_t2_student_example():
    s = slide()
    header(s, "T2 · 시뮬 입력 예시", "가상 학생 1000명 중 1명을 어떻게 설정하는가", 24)

    text(s, 0.7, 2.2, 12, 0.4,
         "1000명을 모두 똑같이 만들면 의미 없음. 학생마다 인지 능력이 다르게 무작위 배정함.",
         size=12, color=MUTED)

    # 좌: 학생 #1 예시 카드
    card(s, 0.7, 2.85, 6.0, 4.3, fill=RGBColor(0xEE, 0xEA, 0xFE), border=ACCENT)
    text(s, 0.95, 3.0, 5.5, 0.4, "예시 학생 #1 · '김민준'", size=14, bold=True, color=ACCENT)
    text(s, 0.95, 3.4, 5.5, 0.35, "코드가 1000명을 만들 때 첫 번째 학생", size=11, color=MUTED)

    line(s, 0.95, 3.85, 6.5, 3.85, color=ACCENT)

    attrs = [
        ("학습률 α", "0.18", "평균(0.15)보다 약간 높음 — 새 경험을 빨리 받아들임"),
        ("망각 d", "0.42", "평균(0.50)보다 낮음 — 기억이 천천히 사라짐"),
        ("작업 기억 WM", "8 chunks", "평균(7)보다 1 큼 — 한 번에 더 많이 처리"),
        ("동기 motivation", "0.67", "평균(0.50)보다 높음 — 자발적 학습 의지 있음"),
    ]
    for i, (name, value, desc) in enumerate(attrs):
        y = 4.0 + i * 0.75
        text(s, 0.95, y, 2.5, 0.35, name, size=11, bold=True, color=INK)
        text(s, 3.4, y, 1.5, 0.4, value, size=18, bold=True, color=ACCENT)
        text(s, 0.95, y + 0.4, 5.5, 0.3, desc, size=11, color=INK_2)

    # 우: 코드 + 해석
    text(s, 7.0, 2.85, 6, 0.4, "어떻게 무작위 추출하는가 (실제 코드)",
         size=11, bold=True, color=ACCENT)

    text(s, 7.0, 3.3, 6, 1.5,
         "rng = np.random.default_rng(42)\n"
         "alpha = rng.normal(0.15, 0.05)\n"
         "d     = rng.normal(0.50, 0.10)\n"
         "WM    = rng.choice([5,6,7,8,9])\n"
         "motiv = rng.normal(0.50, 0.20)",
         size=13, color=INK_2, line_spacing=1.4)

    line(s, 7.0, 4.95, 12.6, 4.95)

    text(s, 7.0, 5.1, 6, 0.4, "1000명 각자 다른 4가지 속성 조합",
         size=11, bold=True, color=ACCENT_2)
    text(s, 7.0, 5.5, 6, 1.6,
         "김민준 같은 학생은 잘 적응한다.\n"
         "반면 α=0.08, WM=5인 학생은 같은 강의계획서에서 훨씬 큰 어려움을 겪는다.\n"
         "1000명 중 가장 약한 25%가 보호되는지를 본다.",
         size=13, color=INK_2, line_spacing=1.4)

    text(s, 0.7, 7.25, 12, 0.3,
         "이 설정값은 문헌 표준값 · T1 설문으로 본 강 학생 실제 값을 수집해 보정 예정",
         size=11, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================
# T2 시뮬 결과 예시 — 분포 효과 (자세히)
# ============================================
def s_t2_examples():
    s = slide()
    header(s, "T2 · 시뮬 결과 예시", "평균만 보면 가려지는 진실 — 분포 효과 차트", 25)

    text(s, 0.7, 2.2, 12, 0.4,
         "같은 강의계획서를 받아도 학생마다 결과가 다르다. 평균보다 \"가장 약한 학생\"에서 차이가 크다.",
         size=14, color=INK_2)

    # 차트는 약간 작게
    chart_path = CHARTS / "chart_distribution.png"
    if chart_path.exists():
        s.shapes.add_picture(str(chart_path), Inches(0.7), Inches(2.85),
                             width=Inches(7.0))

    # 우측 — 자세한 해석 (스토리)
    card(s, 7.95, 2.85, 4.7, 4.2, fill=RGBColor(0xFA, 0xFB, 0xFC), border=ACCENT, border_width=0)

    # 학생 비교: 김민준 vs 홍길동
    text(s, 8.15, 3.0, 4.4, 0.35, "두 학생을 따라가 봅시다", size=11, bold=True, color=ACCENT)

    # 김민준 박스
    card(s, 8.15, 3.45, 4.3, 1.55, fill=RGBColor(0xEC, 0xFD, 0xF5), border=ACCENT_3)
    text(s, 8.35, 3.55, 4, 0.35, "🟢 김민준 (평균 학생)", size=11, bold=True, color=ACCENT_3)
    text(s, 8.35, 3.9, 4, 1.0,
         "기존 강의계획서: 12.5초만에 답 떠올림\n"
         "분산 평가 도입: 10.2초 (조금 빨라짐)",
         size=12, color=INK_2, line_spacing=1.4)

    # 홍길동 박스
    card(s, 8.15, 5.1, 4.3, 1.55, fill=RGBColor(0xFD, 0xE8, 0xF5), border=ACCENT_2)
    text(s, 8.35, 5.2, 4, 0.35, "🔴 홍길동 (약한 학생)", size=11, bold=True, color=ACCENT_2)
    text(s, 8.35, 5.55, 4, 1.0,
         "기존 강의계획서: 22.3초 (1.8배 손해)\n"
         "분산 평가 도입: 13.5초 (대폭 개선)",
         size=12, color=INK_2, line_spacing=1.4)

    text(s, 8.15, 6.75, 4.4, 0.35,
         "→ 같은 개선안인데 효과는 약한 학생에게 더 큼",
         size=10, bold=True, color=INK)

    text(s, 0.7, 7.15, 12, 0.4,
         "결국 우리가 보는 것: 강의계획서가 누구를 가장 힘들게 하는가, 개선안이 그걸 줄이는가",
         size=11, bold=True, color=ACCENT_2, align=PP_ALIGN.CENTER)


# ============================================
# NEW: 개선안 3가지 (이론에서 자동 도출)
# ============================================
def s_improvements():
    s = slide()
    header(s, "OUR PROPOSAL", "우리가 제안하는 개선안 3가지", 20)

    text(s, 0.7, 2.25, 12, 0.4,
         "강의계획서의 문제를 이론으로 분석하면, 개선안이 자연스럽게 따라옵니다.",
         size=13, color=INK_2)

    items = [
        ("01", "분산 평가", "격주 퀴즈 6회 + 기말 40%",
         "기말 한 번이 아니라\n격주마다 작은 퀴즈로 나누어 본다.",
         "Forgetting Curve", ACCENT_4),
        ("02", "프로필 카드", "후보별 정보를 카드로 시각화",
         "팀장이 후보 평가할 때\n역량·시간·관심을 한 화면에 보여준다.",
         "Working Memory 7±2", ACCENT),
        ("03", "명확한 가산", "모호한 감점 → +10 가산만",
         "기준이 모호한 \"감점\" 표현 대신\n명확한 +10 가산만 사용한다.",
         "Framing Effect", ACCENT_2),
    ]
    for i, (num, name, what, mech, theory, color) in enumerate(items):
        x = 0.7 + i * 4.1
        y0 = 2.85
        # 카드 전체 (외곽)
        card(s, x, y0, 3.9, 4.4, fill=WHITE, border=color, border_width=1.5, shadow=True)

        # 상단 컬러 헤더 (카드 안쪽)
        header_h = 1.05
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y0),
                                  Inches(3.9), Inches(header_h))
        bar.fill.solid(); bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        bar.adjustments[0] = 0.05
        # 하단 모서리를 가리는 사각형으로 카드 위만 둥글게
        cover = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(y0 + 0.6),
                                    Inches(3.9), Inches(header_h - 0.6))
        cover.fill.solid(); cover.fill.fore_color.rgb = color
        cover.line.fill.background()

        # 번호 + 이름 (헤더 안)
        text(s, x + 0.3, y0 + 0.2, 1.0, 0.4, num, size=14, bold=True, color=WHITE)
        text(s, x + 0.3, y0 + 0.45, 3.4, 0.55, name, size=22, bold=True, color=WHITE)

        # 부제
        text(s, x + 0.3, y0 + 1.25, 3.4, 0.45, what, size=13, bold=True, color=color)
        # 본문
        text(s, x + 0.3, y0 + 1.85, 3.4, 1.4, mech, size=14, color=INK_2, line_spacing=1.55)

        # 하단 근거 이론
        line(s, x + 0.3, y0 + 3.6, x + 3.6, y0 + 3.6, color=color)
        text(s, x + 0.3, y0 + 3.75, 1.0, 0.35, "근거 이론", size=11, bold=True, color=MUTED)
        text(s, x + 0.3, y0 + 4.0, 3.4, 0.35, theory, size=13, bold=True, color=color)

    text(s, 0.7, 7.4, 12, 0.3,
         "T1 설문은 인식 변화, T2 시뮬은 약한 학생 보호 효과를 검증",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)


# ============================================
# NEW: T3 친분 그래프 시각화 예시
# ============================================
def s_t3_examples():
    s = slide()
    header(s, "T3 · 친분 그래프 예시", "5가지 선택 방식이 만들어내는 그림이 진짜 다름", 24)

    chart_path = CHARTS / "chart_network.png"
    if chart_path.exists():
        s.shapes.add_picture(str(chart_path), Inches(0.7), Inches(2.4),
                             width=Inches(12.0))

    text(s, 0.7, 6.3, 12, 0.4, "한눈에 보이는 차이", size=12, bold=True, color=ACCENT)
    text(s, 0.7, 6.7, 12, 0.4,
         "왼쪽 (친분 우선): 팀장 주변 노드만 분홍색. 친한 사람만 뽑힘.",
         size=12, color=INK_2)
    text(s, 0.7, 7.05, 12, 0.4,
         "오른쪽 (프로필 카드): 그래프 전체에 골고루 분포. 다양한 역량 확보.",
         size=12, color=INK_2)


# ============================================
# 14. 시뮬 결과 시각화 (3장 미니 차트) — KEEP for reference
# ============================================
def s_evidence_grid():
    s = slide()
    header(s, "EVIDENCE", "시뮬레이션이 이미 보여주는 정량 예측치", 14)

    text(s, 0.7, 2.3, 12, 0.4,
         "이 그래프들이 실제 학생 응답에서 동일하게 나오는지 검증하는 것이 우리 미션",
         size=12, color=MUTED)

    chart_paths = [
        (CHARTS / "chart_attendance.png", "출석 정책별 학기 출석률"),
        (CHARTS / "chart_memory.png", "평가 체계별 인출 시간"),
        (CHARTS / "chart_team.png", "후보 수와 평가 완성도"),
    ]
    for i, (path, caption) in enumerate(chart_paths):
        x = 0.7 + i * 4.1
        if path.exists():
            s.shapes.add_picture(str(path), Inches(x), Inches(2.95), width=Inches(3.9))
        text(s, x, 6.4, 3.9, 0.4, caption, size=11, bold=True, color=INK_2, align=PP_ALIGN.CENTER)


# ============================================
# 15. 우리가 가진 자산
# ============================================
def s_assets():
    s = slide()
    header(s, "ASSETS", "우리는 이미 절반 가지고 있다", 15)

    items = [
        ("개인과제 PDF", "8쪽", "이론 매핑 + 시뮬 + 개선안 완성"),
        ("시뮬 코드 3종", "GitHub", "Utility · Memory · WM 모델 구현"),
        ("Streamlit MBTI", "Live", "16유형 게임형 설문 라이브 배포"),
        ("Google Sheets 연동", "실시간", "응답 자동 적재 + 분석 탭 자동"),
        ("Forms 대조군", "28문항", "A/B 비교용 표준 폼 자동 생성"),
        ("Notion 워크스페이스", "5 sub-pages", "설명서·자료실·작업실·회의록·소통방"),
    ]
    for i, (title, badge, desc) in enumerate(items):
        x = 0.7 + (i % 3) * 4.1
        y = 2.3 + (i // 3) * 2.3
        card(s, x, y, 3.9, 2.0)
        text(s, x + 0.3, y + 0.3, 3.5, 0.45, title, size=15, bold=True, color=INK)
        chip(s, x + 0.3, y + 0.85, 1.5, 0.4, badge, color=ACCENT)
        text(s, x + 0.3, y + 1.4, 3.5, 0.55, desc, size=13, color=INK_2, line_spacing=1.3)

    text(s, 0.7, 6.95, 12, 0.4,
         "→ 다른 팀이 처음부터 만들 시간을 우리는 검증·분석에 쓴다.",
         size=12, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)


# ============================================
# 16. 최종 산출물 미리보기
# ============================================
def s_deliverables():
    s = slide()
    header(s, "DELIVERABLES", "보고서·발표 산출물 청사진", 16)

    items = [
        ("📄 보고서 10쪽", "1. 문제 재정의\n2. 이론 매핑 (4개)\n3. 설계 대안 비교\n4. 검증 결과 (T1+T2+T3)\n5. 한계 + 결론"),
        ("🖼️ 핵심 그림 5장", "Fig 1. 4규칙 → 행동 인과도\nFig 2. 출석 utility 곡선\nFig 3. 인출 시간 비교\nFig 4. 팀 구성 완성도 곡선\nFig 5. 친분 네트워크 5조건"),
        ("📊 데이터 표 3개", "표 1. 가설 검정 결과\n표 2. 도구별 응답률 비교\n표 3. 시뮬 vs 실측 정합성"),
        ("🎤 발표 15분", "10분 발표 + 5분 Q&A\n· Streamlit 라이브 데모\n· 네트워크 그래프 시각화\n· 정량 예측 vs 실측 비교"),
    ]
    for i, (title, content) in enumerate(items):
        x = 0.7 + (i % 2) * 6.15
        y = 2.3 + (i // 2) * 2.4
        card(s, x, y, 5.95, 2.15)
        text(s, x + 0.3, y + 0.25, 5.5, 0.5, title, size=15, bold=True, color=ACCENT)
        text(s, x + 0.3, y + 0.85, 5.5, 1.3, content, size=13, color=INK_2, line_spacing=1.4)


# ============================================
# 17. 일정
# ============================================
def s_schedule():
    s = slide()
    header(s, "TIMELINE", "오늘부터 D-32, 마감 6/12", 17)

    weeks = [
        ("11주 · 5/6–5/12", "킥오프 · 주제 확정 · 역할 배분", "TODAY", ACCENT_3),
        ("12주 · 5/13–5/19", "설문 초안 · 시뮬 시나리오 확정 · pilot 5명", "다음 주", INK_2),
        ("13주 · 5/20–5/26", "설문 본 배포 · 시뮬 확장 실행 · 응답 매일 점검", "수집", INK_2),
        ("14주 · 5/28–6/3", "응답 마감 N=60 · 통계 분석 · 결과표", "분석", INK_2),
        ("15주 · 6/4–6/9", "보고서 초고 · 발표자료 · 모의 발표", "마무리", INK_2),
        ("14–15주 수업시간", "최종 발표 (10분 + Q&A 5분)", "발표", ACCENT_2),
        ("16주 · 6/12 23:59", "LMS 최종 보고서 제출", "마감", ACCENT_2),
    ]
    for i, (week, task, status, color) in enumerate(weeks):
        y = 2.3 + i * 0.65
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), Inches(y + 0.18), Inches(0.2), Inches(0.2))
        dot.fill.solid(); dot.fill.fore_color.rgb = color; dot.line.fill.background()
        if i < len(weeks) - 1:
            line(s, 1.1, y + 0.4, 1.1, y + 0.85, color=DIVIDER, weight=1.5)
        text(s, 1.5, y + 0.1, 3.5, 0.4, week, size=12, bold=True, color=INK)
        text(s, 5.2, y + 0.1, 6.5, 0.4, task, size=14, color=INK_2)
        chip(s, 11.4, y + 0.15, 1.4, 0.35, status, color=color)


# ============================================
# 18. 역할 분담
# ============================================
def s_roles():
    s = slide()
    header(s, "DECIDE TODAY", "역할 분담 — 오늘 정한다", 18)

    roles = [
        ("이론·모델링 리드", "이승현", "이미 진행 중", ACCENT_3),
        ("설문 설계·배포", "?", "Forms·Streamlit 운영, 단톡 배포 주도", None),
        ("데이터 분석", "?", "응답 통계 분석 (코드 템플릿 있음)", None),
        ("시뮬·그래프 확장", "?", "T2/T3 코드 실행 및 결과 해석", None),
        ("디자인·발표자료", "?", "PPT 제작 (코드 불필요)", None),
        ("보고서 종합", "?", "10쪽 보고서 총괄", None),
    ]
    for i, (role, person, desc, color) in enumerate(roles):
        y = 2.3 + i * 0.78
        card(s, 0.7, y, 11.9, 0.7, fill=PAPER if color else WHITE)
        text(s, 0.95, y + 0.18, 3.5, 0.4, role, size=14, bold=True, color=INK)
        if color:
            text(s, 4.7, y + 0.15, 1.6, 0.45, person, size=18, bold=True, color=color)
        else:
            text(s, 4.7, y + 0.18, 1.6, 0.45, person, size=18, bold=True, color=MUTED)
        text(s, 6.6, y + 0.2, 6.0, 0.4, desc, size=13, color=INK_2)


# ============================================
# 19. 오늘 결정 4가지
# ============================================
def s_today():
    s = slide()
    header(s, "ACTION", "오늘 끝내야 하는 것 4가지", 19)
    items = [
        ("01", "역할 분담", "위 6개 역할을 누가 어떻게 나눌지 합의"),
        ("02", "정기 회의 시간", "When2meet 결과를 보고 주 1회 시간 확정"),
        ("03", "GitHub 아이디", "팀원 5명 GitHub 아이디 공유 → 저장소 초대"),
        ("04", "설문 배포 시점", "5/13~5/19 중 단톡 배포 일자 합의"),
    ]
    for i, (num, title, desc) in enumerate(items):
        y = 2.4 + i * 1.05
        text(s, 0.95, y + 0.05, 1, 0.7, num, size=36, bold=True, color=ACCENT)
        text(s, 2.0, y, 10, 0.5, title, size=18, bold=True, color=INK)
        text(s, 2.0, y + 0.5, 10, 0.45, desc, size=14, color=INK_2)
        if i < len(items) - 1:
            line(s, 0.7, y + 0.95, 12.6, y + 0.95)
    card(s, 0.7, 6.6, 11.9, 0.6, fill=PAPER)
    text(s, 0.95, 6.72, 11.5, 0.4,
         "+ 이번 주 안에: 각자 개인과제 PDF 자료실 업로드",
         size=11, bold=True, color=ACCENT_2)


# ============================================
# 20. 클로징
# ============================================
def s_closing():
    s = slide(bg=DARK)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(2.4), Inches(0.6), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_2; bar.line.fill.background()

    text(s, 0.7, 2.6, 12, 0.4, "FINAL THOUGHT", size=11, bold=True, color=ACCENT_2)
    text(s, 0.7, 3.05, 12, 1.2, "이미 절반은 되어 있다.",
         size=42, bold=True, color=WHITE)
    text(s, 0.7, 4.2, 12, 0.7, "남은 한 달, 우리는 검증에만 집중한다.",
         size=24, color=MUTED, line_spacing=1.3)
    line(s, 0.7, 6.95, 12.6, 6.95, color=RGBColor(0x35, 0x3B, 0x4D))
    text(s, 0.7, 7.05, 12, 0.4, "Q & A", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─── 슬라이드 추가 ─────────────────────────────
s_cover()                         # 1
s_question()                      # 2
s_4_rules()                       # 3
s_predicted_behavior()            # 4
s_theories()                      # 5
s_theory_utility_intuition()      # 6
s_theory_utility()                # 7
s_theory_forget_intuition()       # 8
s_theory_forget()                 # 9
s_theory_wm_intuition()           # 10
s_theory_wm()                     # 11
s_theory_framing_intuition()      # 12
s_theory_framing()                # 13
s_t1_overview()                   # 14
s_hypotheses()                    # 15
s_t1_attend()                     # 16
s_t1_framing()                    # 17
s_t1_team()                       # 18
s_t1_eval()                       # 19
s_improvements()                  # 20 NEW · 개선안 3가지
s_streamlit_demo_1()              # 21
s_streamlit_demo_2()              # 22
s_t2()                            # 23
s_t2_student_example()            # 24 · 가상 학생 1명 예시
s_t2_examples()                   # 25 · 분포 효과 차트
s_t3()                            # 26
s_t3_examples()                   # 27

out = ROOT / "kickoff_2026-05-11.pptx"
prs.save(out)
print(f"✅ {out}")
print(f"   슬라이드: {len(prs.slides)}장")
